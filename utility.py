import os
import json
import pdfplumber
from google.cloud import storage, aiplatform
import vertexai
from vertexai.language_models import TextEmbeddingModel

def extract_text_from_pdf(pdf_path):
    """Extract all text from a PDF file using pdfplumber."""
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"Error extracting text from {pdf_path}: {e}")
    return text

def extract_text_from_docx(docx_path):
    """Extract all text from a DOCX file using python-docx."""
    from docx import Document
    doc = Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_path):
    """Extract all text from a PPTX file using python-pptx."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_txt(txt_path):
    """Extract all text from a TXT file."""
    with open(txt_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_file(file_path):
    """Dispatch to the correct extractor based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def chunk_text(text, chunk_size=500, overlap=50):
    """Split text into chunks of specified word length with overlap."""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk:
            chunks.append(chunk)
    return chunks

def get_embedding(text, project_id, location):
    """Get embedding vector for a given text using Vertex AI."""
    vertexai.init(project=project_id, location=location)
    model = TextEmbeddingModel.from_pretrained("gemini-embedding-001")
    embeddings = model.get_embeddings([text])
    return embeddings[0].values

def upsert_vector_to_index(chunk_id, embedding, metadata, index_id):
    """Upsert a chunk's embedding and metadata to Vertex AI Matching Engine."""
    embedding = list(embedding)
    restricts = [
        {"namespace": k, "allow_list": [str(v)]}
        for k, v in metadata.items()
    ]
    datapoint = {
        "datapoint_id": chunk_id,
        "feature_vector": embedding,
        "restricts": restricts
    }
    index = aiplatform.MatchingEngineIndex(index_name=index_id)
    index.upsert_datapoints([datapoint])
    print(f"Upserted chunk {chunk_id} to Vertex AI Vector Search.")

def save_chunk_text(chunk_id, chunk_text, chunk_texts, CHUNK_TEXTS_FILE, file_lock=None):
    """
    Save a chunk's text to the chunk_texts mapping and persist to file.
    If file_lock is provided, it is used for thread safety.
    """
    if file_lock:
        file_lock.acquire()
    try:
        chunk_texts[chunk_id] = chunk_text
        temp_file = CHUNK_TEXTS_FILE + ".tmp"
        with open(temp_file, "w", encoding="utf-8") as f:
            json.dump(chunk_texts, f, ensure_ascii=False, indent=2)
        os.replace(temp_file, CHUNK_TEXTS_FILE)
    finally:
        if file_lock:
            file_lock.release()

def download_blob(bucket_name, file_name, destination_folder="downloaded_docs"):
    """Download a file from GCS to a local folder."""
    client = storage.Client()
    bucket = client.bucket(bucket_name)
    blob = bucket.blob(file_name)
    local_path = os.path.join(destination_folder, file_name)
    os.makedirs(os.path.dirname(local_path), exist_ok=True)
    blob.download_to_filename(local_path)
    print(f"Downloaded {file_name} to {local_path}")
    return local_path

def delete_local_file(local_path):
    """Delete a local file if it exists."""
    try:
        os.remove(local_path)
        print(f"Deleted local file: {local_path}")
    except Exception as e:
        print(f"Could not delete local file: {local_path} ({e})")
